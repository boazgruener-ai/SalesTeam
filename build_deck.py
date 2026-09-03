# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- palette: lighter, LinkedIn-brand-toned ----
LI_BLUE = RGBColor(0x0A, 0x66, 0xC2)       # LinkedIn brand blue - headings, accent fills
BLUE = RGBColor(0x5B, 0x9B, 0xD5)          # lighter secondary blue (funnel, minor accents)
LIGHT_BLUE = RGBColor(0xEA, 0xF1, 0xF8)    # soft tint background for callouts
ORANGE = RGBColor(0xD9, 0x73, 0x1E)
DARK = RGBColor(0x22, 0x24, 0x27)
GRAY = RGBColor(0x5B, 0x60, 0x66)
LIGHT_GRAY = RGBColor(0xF4, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1E, 0x7A, 0x4C)
RED = RGBColor(0xA3, 0x2A, 0x2A)

FONT = "Calibri"
PRODUCT_NAME = "SalesTeam"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def textbox(slide, l, t, w, h, text, size=14, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, font=FONT, line_spacing=1.0, anchor=None,
            italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def bullets(slide, l, t, w, h, items, size=15, color=DARK, space_after=10,
            marker="\u2022  ", line_spacing=1.08):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = marker + lead
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.name = FONT
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.name = FONT
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = marker + item
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.color.rgb = color
    return tb


def eyebrow(slide, text):
    textbox(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.35),
            text.upper(), size=12.5, bold=True, color=ORANGE)


def title(slide, text, size=28, top=0.72, width=12.1):
    textbox(slide, Inches(0.6), Inches(top), Inches(width), Inches(0.7),
            text, size=size, bold=True, color=LI_BLUE)


def rule(slide, top=1.45):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(top),
                                   Inches(1.1), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()


def page_footer(slide, n):
    textbox(slide, Inches(12.6), Inches(7.12), Inches(0.6), Inches(0.3),
            str(n), size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    textbox(slide, Inches(0.6), Inches(7.12), Inches(4), Inches(0.3),
            PRODUCT_NAME, size=10, color=GRAY)


def card(slide, l, t, w, h, fill=LIGHT_GRAY, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def make_circle(pic):
    """Mask a picture shape to an ellipse/circle via raw XML (python-pptx has no direct API for this)."""
    spPr = pic._element.spPr
    geom = spPr.find(qn('a:prstGeom'))
    if geom is None:
        geom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'ellipse'})
        geom.append(spPr.makeelement(qn('a:avLst'), {}))
        spPr.append(geom)
    else:
        geom.set('prst', 'ellipse')


PERSONA_PHOTO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "persona-photo.jpg")


def base_slide(n, eyebrow_text, title_text, title_size=28):
    s = add_slide()
    set_bg(s)
    eyebrow(s, eyebrow_text)
    title(s, title_text, size=title_size)
    rule(s)
    page_footer(s, n)
    return s


# ============================================================ SLIDE 1 - TITLE
s = add_slide()
set_bg(s, WHITE)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), SW, Inches(0.06))
band.fill.solid(); band.fill.fore_color.rgb = ORANGE; band.line.fill.background()

textbox(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.3),
        PRODUCT_NAME, size=60, bold=True, color=LI_BLUE)
textbox(s, Inches(0.95), Inches(3.35), Inches(10.5), Inches(0.7),
        "Your AI sales team, right inside LinkedIn.", size=22, color=DARK, italic=True)
textbox(s, Inches(0.95), Inches(4.85), Inches(10.5), Inches(0.5),
        "Backed by a growing team of AI agents \u2014 a Sales Assistant, a Sales Mentor, and a Customer Voice",
        size=14, color=GRAY)
textbox(s, Inches(0.95), Inches(6.65), Inches(9), Inches(0.4),
        "28 August 2026  \u00b7  Presented by Boaz Gruener  \u00b7  AI Product Management Course",
        size=12, color=GRAY)

# ============================================================ SLIDE 2 - PERSONA
s = base_slide(2, "Who It's For", "Meet Elena, B2B Software & Services Sales Manager", title_size=25)
textbox(s, Inches(0.6), Inches(1.65), Inches(9.5), Inches(0.75),
        "Sells enterprise software & AI-driven services in a competitive B2B market. Hits quota by finding "
        "companies actively signaling a need \u2014 not by cold-calling blind.",
        size=13.5, color=GRAY, italic=True, line_spacing=1.15)

if os.path.exists(PERSONA_PHOTO):
    photo = s.shapes.add_picture(PERSONA_PHOTO, Inches(10.55), Inches(1.55), Inches(2.15), Inches(2.15))
    make_circle(photo)
    photo.line.color.rgb = WHITE
    photo.line.width = Pt(3)

quads = [
    ("SAYS", ["\u201cI know the leads are out there somewhere on LinkedIn.\u201d",
              "\u201cI don't have hours to scroll every morning.\u201d"], LIGHT_BLUE, LI_BLUE),
    ("THINKS", ["\u201cDid I miss the perfect post because I searched the wrong week?\u201d",
                "\u201cSales Navigator finds people \u2014 not the moment they're ready to buy.\u201d"], LIGHT_GRAY, DARK),
    ("DOES", ["Runs the same 5 keyword searches by hand, morning after morning.",
              "Skims hundreds of results by eye, one tab at a time."], LIGHT_GRAY, DARK),
    ("FEELS", ["Behind, every day it isn't done.",
               "Relieved on the rare day she finds a great lead early."], LIGHT_BLUE, LI_BLUE),
]
x0, y0, w, h, gap = Inches(0.6), Inches(2.45), Inches(5.85), Inches(2.1), Inches(0.3)
for i, (label, lines, fill, tcol) in enumerate(quads):
    col = i % 2
    row = i // 2
    l = x0 + col * (w + gap)
    t = y0 + row * (h + gap)
    card(s, l, t, w, h, fill=fill)
    textbox(s, l + Inches(0.25), t + Inches(0.18), w - Inches(0.5), Inches(0.3), label,
            size=13, bold=True, color=ORANGE)
    bullets(s, l + Inches(0.25), t + Inches(0.55), w - Inches(0.5), h - Inches(0.7),
            lines, size=12.5, color=tcol, space_after=8)

# ============================================================ SLIDE 3 - PROBLEM & WHY NOW (merged)
s = base_slide(3, "The Problem & Why Now", "Lead generation eats the day \u2014 and the gap is widening", title_size=25)
colw, colgap = Inches(5.65), Inches(0.3)
colx = [Inches(0.7), Inches(0.7) + colw + colgap]

textbox(s, colx[0], Inches(1.75), colw, Inches(0.35), "THE PROBLEM", size=14, bold=True, color=ORANGE)
bullets(s, colx[0], Inches(2.15), colw, Inches(3.3), [
    ("Repetitive manual work. ", "Finding buying-intent signals means running the same keyword "
     "searches by hand, across Posts and Jobs, every day."),
    ("Wrong tool, wrong job. ", "Sales Navigator and enrichment tools find contact info for a name "
     "you already have \u2014 not who's signaling a need right now."),
    ("The cost. ", "Hours lost weekly to manual searching \u2014 or leads simply missed."),
], size=13.5, space_after=16)

textbox(s, colx[1], Inches(1.75), colw, Inches(0.35), "WHY NOW", size=14, bold=True, color=ORANGE)
bullets(s, colx[1], Inches(2.15), colw, Inches(3.3), [
    ("More noise. ", "LinkedIn is bigger and busier than ever \u2014 more signal buried under more "
     "scroll."),
    ("Riskier shortcuts. ", "Automation tools like Waalaxy, Expandi and Dripify saw a wave of account "
     "restrictions in Jan\u2013Mar 2026."),
    ("The platform caught up. ", "Manifest V3, Chrome's side panel, and Claude's tool-use API make a "
     "careful, reasoning-backed assistant possible now."),
], size=13.5, space_after=16)

card(s, Inches(0.7), Inches(5.7), Inches(11.6), Inches(1.0), fill=LIGHT_BLUE)
textbox(s, Inches(1.0), Inches(5.85), Inches(11), Inches(0.7),
        "\u201cGenerate sales leads on LinkedIn without spending hours on repetitive manual searches \u2014 "
        "and without investing in big, complicated sales tools.\u201d", size=14.5, italic=True, color=LI_BLUE)

# ============================================================ SLIDE 4 - SOLUTION
s = base_slide(4, "The Solution", "Not just a leads generator \u2014 an AI sales team in your browser", title_size=25)
card(s, Inches(0.7), Inches(1.7), Inches(11.6), Inches(0.85), fill=LI_BLUE)
textbox(s, Inches(1.0), Inches(1.85), Inches(11), Inches(0.6),
        "An easy-to-use browser-embedded tool that performs smart, automated searches with "
        "custom targeting via keywords and filters.", size=15, italic=True, color=WHITE)
bullets(s, Inches(0.7), Inches(2.8), Inches(11.6), Inches(3.5), [
    ("Topics. ", "Define reusable keyword groups, with an optional second \u201cAND with\u201d condition."),
    ("One click. ", "\u201cScan All Topics\u201d searches every Topic across both LinkedIn Posts and Jobs."),
    ("One ranked list. ", "Results are merged, deduplicated and ranked \u2014 tagged with matched "
     "keywords, hiring/freelance signals, and job-ad detection. Competitors and recruiter noise are "
     "auto-filtered out, not just left for you to skip past."),
    ("Safe by design. ", "100% manual-trigger. Nothing runs unattended, so it stays on the safe side "
     "of LinkedIn's policies."),
    ("An AI sales team on top. ", "Specialist agents that prioritize leads, draft outreach, and pressure-test "
     "your message \u2014 detailed next."),
], size=14.5, space_after=9)

# ============================================================ SLIDE 5 - THE AI SALES TEAM
s = base_slide(5, "The AI Sales Team", "A growing team of specialist agents, on call 24/7", title_size=25)
w = Inches(3.75); gap = Inches(0.25)
team = [
    ("SALES ASSISTANT", "Full pipeline agent \u2014 lead operations.", [
        "Scans Posts + Jobs, auto-filters out competitors and recruiter noise, and tracks every lead's status "
        "end to end.",
        "Scores each new lead P1\u2013P5 by real fit and urgency after every scan \u2014 so nothing falls "
        "through the cracks.",
    ], LIGHT_GRAY),
    ("SALES MENTOR", "AI Agent \u00b7 25 years of experience.", [
        "Ask which leads to prioritize, how to approach one \u2014 or a general strategy question.",
        "Decides for itself when it needs to check real lead data before answering.",
    ], LIGHT_BLUE),
    ("CUSTOMER VOICE", "AI Agent \u00b7 realistic buyer persona.", [
        "Bounce a sales pitch or message off it before it ever reaches a real lead.",
        "Ask it about needs, pains, or how it prefers to be approached \u2014 an open-ended interview, not a "
        "fixed script.",
    ], LIGHT_GRAY),
]
for i, (head, sub, items, fill) in enumerate(team):
    l = Inches(0.6) + i * (w + gap)
    card(s, l, Inches(1.85), w, Inches(3.3), fill=fill)
    textbox(s, l + Inches(0.25), Inches(2.05), w - Inches(0.5), Inches(0.35), head, size=14, bold=True, color=ORANGE)
    textbox(s, l + Inches(0.25), Inches(2.4), w - Inches(0.5), Inches(0.5), sub, size=11.5, italic=True, color=GRAY)
    bullets(s, l + Inches(0.25), Inches(2.9), w - Inches(0.5), Inches(2.1), items, size=12, space_after=10)
card(s, Inches(0.6), Inches(5.35), Inches(11.6), Inches(1.35), fill=LI_BLUE)
textbox(s, Inches(0.9), Inches(5.5), Inches(11), Inches(1.05),
        "One shared tool-use engine powers every agent on the team \u2014 each just gets a different persona "
        "and a different set of tools it can call on its own. That's genuine agentic AI: the model decides "
        "what it needs to look up, not a scripted chatbot. More specialist agents join the team as new needs "
        "surface.",
        size=13, italic=True, color=WHITE, line_spacing=1.2)

# ============================================================ SLIDE 6 - LIVE DEMO
s = base_slide(6, "Live Demo", "See it live \u2192")
steps = [
    ("1", "Define Topics", "What to find (e.g. \u201cAI Transformation\u201d) and what to auto-filter out (competitors, recruiters)."),
    ("2", "Click \u201cScan All Topics\u201d", "Searches Posts + Jobs, filters the noise, then scores every new lead P1\u2013P5."),
    ("3", "Open the Dashboard", "One sortable, filterable pipeline \u2014 highest-priority leads on top, low-priority hidden by default."),
    ("4", "Ask Sales Mentor for advice", "\u201cWhich lead should I approach first?\u201d \u2014 grounded in your real results."),
    ("5", "Ask Sales Mentor to draft it", "Turns its own recommendation into a ready-to-send opening line."),
]
w, gap = Inches(2.2), Inches(0.25)
x0 = Inches(0.6)
for i, (num, head, body) in enumerate(steps):
    l = x0 + i * (w + gap)
    card(s, l, Inches(2.0), w, Inches(3.7), fill=LIGHT_GRAY)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.2), Inches(2.25), Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = ORANGE; circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = WHITE
    textbox(s, l + Inches(0.18), Inches(2.9), w - Inches(0.36), Inches(0.75), head, size=13.5, bold=True,
            color=LI_BLUE, line_spacing=1.05)
    textbox(s, l + Inches(0.18), Inches(3.7), w - Inches(0.36), Inches(1.85), body, size=11, color=GRAY,
            line_spacing=1.15)
textbox(s, Inches(0.7), Inches(5.95), Inches(11), Inches(0.5),
        "(Switch to the extension's side panel now for the live walkthrough.)", size=12, italic=True, color=GRAY)

# ============================================================ SLIDE 7 - HOW IT WORKS + ARCHITECTURE (merged)
s = base_slide(7, "How It Works", "Runs inside your own LinkedIn session \u2014 architecture underneath", title_size=24)
bullets(s, Inches(0.7), Inches(1.5), Inches(11.6), Inches(0.85), [
    ("Native browser extension. ", "Chrome/Edge, Manifest V3, in the side panel next to LinkedIn."),
    ("Adaptive query-chunking. ", "Works around LinkedIn's undocumented search-complexity limits automatically."),
    ("Agents call tools, not scripts. ", "Claude decides, mid-conversation, when to fetch real lead data."),
    ("Your data stays yours. ", "Only the AI features send lead text to Anthropic, using your own key."),
], size=11, space_after=3, line_spacing=1.0)


def arch_box(l, t, w, h, label, sub, fill, tcol):
    shp = card(s, l, t, w, h, fill=fill)
    textbox(s, l, t + Inches(0.1), w, Inches(0.3), label, size=12.5, bold=True, color=tcol, align=PP_ALIGN.CENTER)
    textbox(s, l + Inches(0.1), t + Inches(0.42), w - Inches(0.2), h - Inches(0.5), sub, size=9,
            color=tcol, align=PP_ALIGN.CENTER, line_spacing=1.05)
    return shp


linkedin = arch_box(Inches(0.5), Inches(3.3), Inches(2.0), Inches(1.35), "LinkedIn",
                     "Third-party \u00b7 Posts & Jobs pages", LIGHT_GRAY, GRAY)
side_panel = arch_box(Inches(3.05), Inches(2.85), Inches(3.05), Inches(1.3), "Side Panel UI",
                       "Built \u00b7 Topics, Negative\nTopics, scan trigger", LIGHT_BLUE, LI_BLUE)
bg_worker = arch_box(Inches(6.35), Inches(2.85), Inches(3.05), Inches(1.3), "Background Worker",
                      "Built \u00b7 plans searches,\ndrives the scan tab", LIGHT_BLUE, LI_BLUE)
content_script = arch_box(Inches(6.35), Inches(4.4), Inches(3.05), Inches(1.05), "Content Scripts",
                           "Built \u00b7 reads LinkedIn's\nrendered page", LIGHT_BLUE, LI_BLUE)
storage = arch_box(Inches(3.05), Inches(4.4), Inches(3.05), Inches(1.05), "Local Storage",
                    "Browser API \u00b7 on this device only", LIGHT_GRAY, DARK)
anthropic = arch_box(Inches(9.9), Inches(3.3), Inches(2.85), Inches(1.35), "Anthropic API",
                      "Third-party \u00b7 drafts + agents,\nyour own key", LIGHT_GRAY, GRAY)

# Straight connectors between box edges, each with a short label naming what moves.
# LinkedIn <-> Side Panel/Storage column
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.5), Inches(3.65), Inches(3.05), Inches(3.5)).line.color.rgb = GRAY
textbox(s, Inches(2.35), Inches(3.85), Inches(1.0), Inches(0.5), "navigate /\nscrape", size=8, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.0)

# Side Panel <-> Background Worker
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.1), Inches(3.5), Inches(6.35), Inches(3.5)).line.color.rgb = GRAY
textbox(s, Inches(5.6), Inches(3.58), Inches(1.15), Inches(0.5), "scan\nrequest /\nprogress", size=7.5, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=0.95)

# Content Scripts <-> Background Worker (vertical, same column)
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.87), Inches(4.15), Inches(7.87), Inches(4.4)).line.color.rgb = GRAY
textbox(s, Inches(8.0), Inches(4.18), Inches(1.4), Inches(0.25), "scraped leads", size=8, color=GRAY, align=PP_ALIGN.LEFT)

# Side Panel <-> Local Storage (vertical, same column)
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.57), Inches(4.15), Inches(4.57), Inches(4.4)).line.color.rgb = GRAY
textbox(s, Inches(2.55), Inches(4.18), Inches(1.9), Inches(0.25), "settings + results", size=8, color=GRAY, align=PP_ALIGN.RIGHT)

# Local Storage <-> Content Scripts (background worker also reads/writes storage, mediated by itself above)
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.1), Inches(4.9), Inches(6.35), Inches(4.9)).line.color.rgb = GRAY
textbox(s, Inches(5.6), Inches(4.98), Inches(1.15), Inches(0.4), "reads /\nwrites", size=7.5, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=0.95)

# Side Panel <-> Anthropic API - routed above the row so it doesn't cut through
# Background Worker's box in between.
elbow_y = Inches(2.6)
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.575), Inches(2.85), Inches(4.575), elbow_y).line.color.rgb = GRAY
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.575), elbow_y, Inches(11.325), elbow_y).line.color.rgb = GRAY
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(11.325), elbow_y, Inches(11.325), Inches(3.3)).line.color.rgb = GRAY
textbox(s, Inches(5.6), Inches(2.4), Inches(4.7), Inches(0.2), "prompt + tools \u2192   \u2190 draft / tool call",
        size=8, color=GRAY, align=PP_ALIGN.CENTER)

textbox(s, Inches(0.6), Inches(5.7), Inches(11.6), Inches(0.4),
        "No server of our own \u2014 everything above runs on your device or a third-party API you already trust. "
        "Three more pages now share this same local data: a Dashboard (pipeline + AI prioritization), "
        "Advisors, and Settings.",
        size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================ SLIDE 8 - MARKET OPPORTUNITY
s = base_slide(8, "Market Opportunity", "Grounded in real market data, not invented numbers", title_size=24)

tiers = [
    ("TAM", "$4.85B", "~8M B2B sellers worldwide", "Global sales intelligence software market, 2025",
     "Sourced \u2193", LI_BLUE),
    ("SAM", "~$730M", "~1.2M individual / SMB sellers", "Our estimate: ~15% of TAM \u2014 the slice without a "
     "full enterprise sales-intelligence budget", "Our estimate", BLUE),
    ("SOM", "~$36.5M", "~60K potential buyers", "Our estimate: ~5% of SAM \u2014 realistically reachable over "
     "several years, not a year-1 target", "Our estimate", ORANGE),
]
tl, tw, th, tgap = Inches(0.6), Inches(5.6), Inches(1.35), Inches(0.15)
ty = Inches(1.65)
for label, value, buyers, desc, tag, color in tiers:
    card(s, tl, ty, tw, th, fill=color)
    textbox(s, tl + Inches(0.25), ty + Inches(0.1), tw - Inches(0.5), Inches(0.35),
            f"{label}  \u2014  {value}", size=16, bold=True, color=WHITE)
    textbox(s, tl + Inches(0.25), ty + Inches(0.46), tw - Inches(0.5), Inches(0.25),
            buyers, size=11, bold=True, color=WHITE)
    textbox(s, tl + Inches(0.25), ty + Inches(0.74), tw - Inches(0.5), Inches(0.55),
            f"{desc}  \u00b7  {tag}", size=9, italic=True, color=WHITE, line_spacing=1.05)
    ty = ty + th + tgap

# 5-year TAM growth chart (right column)
cx0, cw = Inches(6.7), Inches(5.9)
textbox(s, cx0, Inches(1.65), cw, Inches(0.35), "TAM growth, 2025\u20132030", size=14, bold=True, color=LI_BLUE)
textbox(s, cx0, Inches(2.0), cw, Inches(0.3), "11.10% CAGR \u00b7 Fortune Business Insights", size=10, italic=True, color=GRAY)

years = ["2025", "2026", "2027", "2028", "2029", "2030"]
values = [4.85, 5.37, 5.97, 6.63, 7.37, 8.19]
baseline = Inches(5.55)
max_h = Inches(2.85)
bar_w, bar_gap = Inches(0.68), Inches(0.16)
n = len(values)
chart_w = n * bar_w + (n - 1) * bar_gap
bx0 = cx0 + (cw - chart_w) / 2
for i, (yr, val) in enumerate(zip(years, values)):
    bh = int(max_h * (val / max(values)))
    bx = bx0 + i * (bar_w + bar_gap)
    by = baseline - bh
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bar_w, Emu(bh))
    bar.adjustments[0] = 0.12
    bar.fill.solid(); bar.fill.fore_color.rgb = LI_BLUE if i == 0 else BLUE
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s, bx - Inches(0.1), by - Inches(0.28), bar_w + Inches(0.2), Inches(0.25),
            f"${val:.2f}B", size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    textbox(s, bx - Inches(0.1), baseline + Inches(0.05), bar_w + Inches(0.2), Inches(0.25),
            yr, size=9.5, color=GRAY, align=PP_ALIGN.CENTER)

card(s, Inches(0.6), Inches(6.2), Inches(11.9), Inches(0.85), fill=LIGHT_GRAY)
textbox(s, Inches(0.85), Inches(6.28), Inches(11.4), Inches(0.72),
        "Sources: TAM and 5-year projection (11.10% CAGR) from Fortune Business Insights' Sales Intelligence "
        "Market report. Buyer counts extrapolated from U.S. Bureau/industry data (~2.5\u20133M B2B sellers in the "
        "U.S.) scaled to a global estimate \u2014 not an independently sourced global figure. SAM and SOM buyer "
        "counts and market size are our own estimates.",
        size=9.5, italic=True, color=GRAY, line_spacing=1.1)

# ============================================================ SLIDE 9 - COMPETITIVE LANDSCAPE
s = base_slide(9, "The Landscape", "Everyone solves a slice. We watch the signal \u2014 and coach the play.", title_size=24)
cols = ["", "Runs in your\nsession", "Finds intent\nsignals", "No outreach-\nautomation risk", "AI sales team\nfor support", "Price"]
rows = [
    ["Sales Navigator", "\u2713", "\u2717", "\u2713", "\u2717", "$99.99/mo/seat"],
    ["Contact enrichment\n(Apollo, Lusha, Kaspr, ZoomInfo)", "\u2717", "\u2717", "\u2713", "\u2717", "$39\u2013166/mo"],
    ["Outreach automation\n(Waalaxy, PhantomBuster, Dux-Soup...)", "\u25CF", "\u25CF", "\u2717", "\u2717", "$30\u2013120/mo"],
    ["Manual browsing", "\u2713", "\u25CF", "\u2713", "\u2717", "Free"],
    [PRODUCT_NAME, "\u2713", "\u2713", "\u2713", "\u2713", "Low-cost"],
]
tb_top = Inches(1.7)
tb_left = Inches(0.5)
tb_w = Inches(12.35)
tb_h = Inches(3.75)
gtable = s.shapes.add_table(len(rows) + 1, len(cols), tb_left, tb_top, tb_w, tb_h).table
gtable.columns[0].width = Inches(3.75)
rest_w = int((tb_w - Inches(3.75)) / (len(cols) - 1))
for i in range(1, len(cols)):
    gtable.columns[i].width = Emu(rest_w)

for j, colname in enumerate(cols):
    cell = gtable.cell(0, j)
    cell.text = colname
    cell.fill.solid(); cell.fill.fore_color.rgb = LI_BLUE
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT

for i, row in enumerate(rows, start=1):
    is_us = row[0] == PRODUCT_NAME
    for j, val in enumerate(row):
        cell = gtable.cell(i, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE if is_us else (WHITE if i % 2 else LIGHT_GRAY)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.name = FONT
                r.font.bold = is_us or j == 0
                if j in (1, 2, 3, 4):
                    r.font.color.rgb = GREEN if val == "\u2713" else (ORANGE if val == "\u25CF" else RED)
                else:
                    r.font.color.rgb = LI_BLUE if is_us else DARK
    gtable.rows[i].height = Inches(0.62)

textbox(s, Inches(0.5), Inches(5.65), Inches(12.35), Inches(0.5),
        "\u25CF = partial   \u2717 = no. ~40% of accounts on flagged automation tools were restricted in a single "
        "quarter (Jan\u2013Mar 2026). No competitor offers an AI sales team backing up the seller.", size=11.5,
        italic=True, color=GRAY)

# ============================================================ SLIDE 10 - EARLY VALIDATION
s = base_slide(10, "Early Validation", "A working alpha, not a mockup")
bullets(s, Inches(0.7), Inches(1.9), Inches(11.6), Inches(3.4), [
    ("Shipped. ", "v0.6 built and running today \u2014 Posts + Jobs scanning (English & German), AI-drafted "
     "outreach, and working AI agents: Sales Mentor and Customer Voice, with the Sales Assistant's own agent "
     "capabilities next."),
    ("Proven at scale, live. ", "After a real optimization pass, one batch scan surfaced 75 ranked leads "
     "across Posts and Jobs \u2014 up from 12 in the first run \u2014 with zero manual searching."),
    ("Real agent conversations, not mockups. ", "Sales Mentor and Customer Voice already hold real "
     "tool-calling conversations grounded in scraped lead data."),
    ("A second real user, real feedback. ", "In active use by a second B2B salesperson \u2014 a real sales "
     "team lead \u2014 whose real usage has already surfaced and fixed genuine reliability bugs."),
], size=14, space_after=14)
card(s, Inches(0.7), Inches(5.75), Inches(11.6), Inches(1.1), fill=LI_BLUE)
textbox(s, Inches(1.0), Inches(5.9), Inches(11), Inches(0.8),
        "This isn't a concept slide \u2014 it's a product (and an AI sales team) that ran on the real LinkedIn "
        "this week.", size=15, italic=True, color=WHITE)

# ============================================================ SLIDE 11 - BUSINESS MODEL & ECONOMICS (merged)
s = base_slide(11, "Business Model & Economics \u00b7 Illustrative", "Freemium pricing, illustrative unit economics", title_size=25)
w = Inches(5.85); gap = Inches(0.3)
tiers = [
    ("FREE", "$0", ["1 active Topic", "Posts search only", "Manual export"], LIGHT_GRAY, LI_BLUE),
    ("PRO", "$12/mo", ["Unlimited Topics + Job Search", "AI drafting, Sales Mentor & Customer Voice",
                        "CSV export + settings sync"], LI_BLUE, WHITE),
]
for i, (name, price, feats, fill, tcol) in enumerate(tiers):
    l = Inches(0.6) + i * (w + gap)
    card(s, l, Inches(1.65), w, Inches(1.95), fill=fill)
    textbox(s, l + Inches(0.3), Inches(1.8), Inches(2.2), Inches(0.35), name, size=14, bold=True,
            color=(ORANGE if i else LI_BLUE))
    textbox(s, l + Inches(0.3), Inches(2.1), Inches(2.2), Inches(0.5), price, size=22, bold=True, color=tcol)
    bullets(s, l + Inches(2.6), Inches(1.85), w - Inches(2.9), Inches(1.65), feats, size=11.5, color=tcol, space_after=6)

metrics = [
    ("50,000", "Year 1 users (target)"),
    ("5%", "Free \u2192 Paid conversion"),
    ("2,500", "Paying users"),
    ("$12", "ARPU / month"),
    ("~$360K", "Year 1 ARR"),
]
mw = Inches(2.25); mgap = Inches(0.15)
for i, (num, label) in enumerate(metrics):
    l = Inches(0.6) + i * (mw + mgap)
    card(s, l, Inches(3.85), mw, Inches(1.35), fill=LIGHT_BLUE if i == len(metrics) - 1 else LIGHT_GRAY)
    textbox(s, l, Inches(3.98), mw, Inches(0.55), num, size=19, bold=True, color=LI_BLUE, align=PP_ALIGN.CENTER)
    textbox(s, l + Inches(0.1), Inches(4.55), mw - Inches(0.2), Inches(0.6), label, size=10.5, color=GRAY,
            align=PP_ALIGN.CENTER, line_spacing=1.05)

card(s, Inches(0.6), Inches(5.45), Inches(11.6), Inches(1.35), fill=LIGHT_GRAY)
textbox(s, Inches(0.9), Inches(5.6), Inches(11), Inches(1.1),
        "North star metric: Weekly Active Scanners. Pricing and economics above are illustrative placeholders "
        "for a course exercise \u2014 not based on real pricing tests, real conversion data, or market "
        "research, and should not be read as a forecast.",
        size=12, italic=True, color=GRAY, line_spacing=1.2)

# ============================================================ SLIDE 12 - ROADMAP + ASK
s = base_slide(12, "Where It Goes \u00b7 The Ask", "Roadmap and what we need next", title_size=26)
stages = [
    ("NOW \u00b7 v0.13", ["Posts + Jobs scanning, English & German, auto-filtering out competitors and "
     "recruiter noise", "A full pipeline Dashboard: every lead tracked, AI-scored P1\u2013P5, safe bulk "
     "actions with undo", "AI-drafted outreach, Sales Mentor & Customer Voice agents"], LIGHT_GRAY),
    ("NEXT", ["Push validated, high-priority leads straight into Salesforce or another CRM", "Integrating "
     "with Sales Navigator to enhance search and filtering capabilities"], LIGHT_GRAY),
    ("LATER", ["Chrome Web Store distribution", "Team workspaces"], LIGHT_GRAY),
]
w = Inches(3.75); gap = Inches(0.25)
for i, (head, items, fill) in enumerate(stages):
    l = Inches(0.6) + i * (w + gap)
    card(s, l, Inches(1.85), w, Inches(2.8), fill=fill)
    textbox(s, l + Inches(0.25), Inches(2.05), w - Inches(0.5), Inches(0.4), head, size=14, bold=True, color=ORANGE)
    bullets(s, l + Inches(0.25), Inches(2.55), w - Inches(0.5), Inches(2.0), items, size=11, space_after=7)
card(s, Inches(0.6), Inches(4.95), Inches(11.6), Inches(1.35), fill=LI_BLUE)
textbox(s, Inches(0.9), Inches(5.1), Inches(0.9), Inches(1.0), "THE\nASK", size=14, bold=True, color=ORANGE, line_spacing=1.0)
textbox(s, Inches(2.0), Inches(5.15), Inches(10), Inches(1.0),
        "Feedback on this direction, and introductions to 3\u20135 pilot sales teams willing to alpha-test the "
        "full AI sales team \u2014 scanning, drafting, and advisors \u2014 against their real pipelines.",
        size=14, color=WHITE, line_spacing=1.2)

# ============================================================ SLIDE 15 - VISION + THANK YOU
s = add_slide()
set_bg(s, WHITE)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), SW, Inches(0.06))
band.fill.solid(); band.fill.fore_color.rgb = ORANGE; band.line.fill.background()
textbox(s, Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.8),
        "Every LinkedIn seller deserves an AI sales team behind them \u2014 not a blank search bar.",
        size=28, bold=True, color=LI_BLUE, line_spacing=1.15)
textbox(s, Inches(1.0), Inches(4.85), Inches(9), Inches(0.5), "Thank you.", size=20, color=DARK)
textbox(s, Inches(1.0), Inches(5.5), Inches(9), Inches(0.5),
        "Boaz Gruener  \u00b7  boaz.gruener@gmail.com", size=13, color=GRAY)

# ============================================================ APPENDIX DIVIDER
s = add_slide()
set_bg(s, LIGHT_GRAY)
textbox(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.0),
        "APPENDIX", size=44, bold=True, color=LI_BLUE)
textbox(s, Inches(0.95), Inches(3.95), Inches(10), Inches(0.5),
        "Backup material \u2014 shown only if it comes up in Q&A.", size=15, italic=True, color=GRAY)

# ============================================================ APPENDIX A1 - TWO SEGMENTS (moved from main deck)
s = base_slide("A1", "Two Segments, One Engine", "Works with or without Sales Navigator")
w = Inches(5.85); gap = Inches(0.3)
segs = [
    ("WITHOUT Sales Navigator", LIGHT_GRAY, [
        "The core use case \u2014 works on any LinkedIn account, free or paid.",
        "Topics become the only way to systematically surface buying-intent signals.",
        "No added subscription cost on top of what the seller already pays (if anything).",
    ]),
    ("WITH Sales Navigator", LIGHT_BLUE, [
        "The same Topics engine layers on top of Sales Navigator's richer people/company filters.",
        "Sharpens targeting instead of replacing the subscription \u2014 makes an existing spend work harder.",
        "Roadmap: direct integration with Sales Navigator's saved searches.",
    ]),
]
for i, (head, fill, items) in enumerate(segs):
    l = Inches(0.6) + i * (w + gap)
    card(s, l, Inches(1.85), w, Inches(4.2), fill=fill)
    textbox(s, l + Inches(0.3), Inches(2.1), w - Inches(0.6), Inches(0.5), head, size=16, bold=True, color=LI_BLUE)
    bullets(s, l + Inches(0.3), Inches(2.7), w - Inches(0.6), Inches(3.1), items, size=13.5, space_after=14)

out_path = "SalesTeam - Pitch Deck.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
